import os
import sys
import shutil
import site

import win32com.client
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import load_workbook


# -----------------------------
# Fix pywin32 DLL issue
# -----------------------------

def fix_pywin32():
    site_packages = site.getsitepackages()[0]
    pywin32_system32 = os.path.join(site_packages, "pywin32_system32")
    python_dlls = os.path.join(sys.base_prefix, "DLLs")

    if os.path.exists(pywin32_system32):
        for f in os.listdir(pywin32_system32):
            if f.lower().startswith(("pywintypes", "pythoncom")):
                src = os.path.join(pywin32_system32, f)
                dst = os.path.join(python_dlls, f)
                if not os.path.exists(dst):
                    shutil.copy(src, dst)


fix_pywin32()


# -----------------------------
# Config
# -----------------------------

excel_file = "Excel Sample.xlsx"

wb = load_workbook(excel_file, data_only=True)

prs = Presentation()


# -----------------------------
# Utility Functions
# -----------------------------

def create_blank_slide():

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.element.getparent().remove(shape.element)

    return slide


def hex_to_rgb(hex_color):

    if not hex_color:
        return None

    hex_color = hex_color[-6:]

    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def style_cell(ppt_cell, excel_cell, font_size=11):

    value = excel_cell.value
    ppt_cell.text = "" if value is None else str(value)

    fill = excel_cell.fill

    if fill and fill.start_color and fill.start_color.rgb not in [None, "00000000", "FFFFFFFF"]:
        rgb = hex_to_rgb(fill.start_color.rgb)

        if rgb:
            ppt_cell.fill.solid()
            ppt_cell.fill.fore_color.rgb = rgb

    for paragraph in ppt_cell.text_frame.paragraphs:
        for run in paragraph.runs:

            run.font.size = Pt(font_size)

            font_color = excel_cell.font.color

            if font_color and font_color.rgb:
                rgb = hex_to_rgb(font_color.rgb)

                if rgb:
                    run.font.color.rgb = rgb

            run.font.bold = excel_cell.font.bold


def add_excel_table(slide, sheet, start_row, end_row, start_col, end_col, top):

    rows = end_row - start_row + 1
    cols = end_col - start_col + 1

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        top,
        Inches(9),
        Inches(3)
    ).table

    for r in range(rows):
        for c in range(cols):

            excel_cell = sheet.cell(row=start_row + r, column=start_col + c)
            ppt_cell = table.cell(r, c)

            style_cell(ppt_cell, excel_cell)


def add_large_table(slide, sheet):

    rows, cols = 5, 2

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1),
        Inches(9),
        Inches(5)
    ).table

    for r in range(rows):
        for c in range(cols):

            excel_cell = sheet.cell(row=r+1, column=c+1)
            ppt_cell = table.cell(r, c)

            style_cell(ppt_cell, excel_cell, font_size=32)


def add_excel_charts(slide, excel_file):

    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False

    wb_excel = excel.Workbooks.Open(os.path.abspath(excel_file))
    sheet = wb_excel.Sheets(2)

    left = Inches(0.5)
    top = Inches(1.5)

    chart_index = 1

    for shape in sheet.Shapes:

        if shape.HasChart:

            chart = shape.Chart

            image_path = os.path.abspath(f"chart_{chart_index}.png")

            chart.Export(image_path)

            slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=Inches(4.5),
                height=Inches(3)
            )

            left += Inches(4.7)

            if left > Inches(8):
                left = Inches(0.5)
                top += Inches(3.2)

            chart_index += 1

    wb_excel.Close(False)
    excel.Quit()


# -----------------------------
# Create Slides
# -----------------------------

slide1 = create_blank_slide()
slide2 = create_blank_slide()
slide3 = create_blank_slide()


# Slide 1
sheet3 = wb["3"]
add_large_table(slide1, sheet3)


# Slide 2
add_excel_charts(slide2, excel_file)


# Slide 3
sheet1 = wb["1"]
add_excel_table(slide3, sheet1, 1, 12, 1, 8, Inches(1))


# -----------------------------
# Save
# -----------------------------

prs.save("Excel_to_PPT_Final_Output.pptx")

print("PowerPoint created successfully.")